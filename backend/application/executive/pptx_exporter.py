"""M48.2 G-018 — Board Report PPTX Export via python-pptx.

Generates a PowerPoint presentation from a BoardReportModel.

Slides:
  1. Cover — org name, report title, date
  2. Executive Summary — key metrics grid
  3. KPI Trends — text summary (charts require a rendering layer)
  4. Risk Register — top 10 open risks
  5. Next Steps — pending recommendations

Architecture:
  - Pure python-pptx; no external rendering service.
  - Returns bytes (caller streams via FastAPI StreamingResponse).
  - No generative AI — all content comes from the DB report model.
"""

from __future__ import annotations

import io
from datetime import date
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

_BLUE = RGBColor(0x1D, 0x4E, 0xD8)    # primary brand blue
_DARK = RGBColor(0x0F, 0x17, 0x2A)    # near-black
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_LIGHT = RGBColor(0xF1, 0xF5, 0xF9)   # slate-100
_ACCENT = RGBColor(0x10, 0xB9, 0x81)  # emerald-500


def _set_cell_bg(cell, color: RGBColor) -> None:
    from pptx.oxml.ns import qn
    from lxml import etree
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", f"{color[0]:02x}{color[1]:02x}{color[2]:02x}")


def _add_slide(prs: Presentation, layout_idx: int = 6):
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def _text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def _fill_background(slide, color: RGBColor) -> None:
    from pptx.oxml.ns import qn
    from lxml import etree
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def build_board_report_pptx(
    *,
    organization_name: str,
    report_title: str,
    report_date: date | str | None = None,
    executive_summary: str = "",
    kpi_highlights: list[dict[str, Any]] | None = None,
    risks: list[dict[str, Any]] | None = None,
    recommendations: list[dict[str, Any]] | None = None,
    period_label: str = "",
) -> bytes:
    """Build a BoardReport PPTX and return it as bytes.

    Args:
        kpi_highlights: list of {label, value, unit, trend}
        risks: list of {title, severity, status, owner}
        recommendations: list of {title, priority, due_date}

    Returns:
        PPTX file as bytes.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    kpi_highlights = kpi_highlights or []
    risks = risks or []
    recommendations = recommendations or []
    date_str = str(report_date) if report_date else str(date.today())

    # ── Slide 1: Cover ────────────────────────────────────────────────────────
    slide = _add_slide(prs)
    _fill_background(slide, _DARK)

    _banner = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), Inches(13.33), Inches(1.5)
    )
    _banner.fill.solid()
    _banner.fill.fore_color.rgb = _BLUE

    _text_box(slide, 0.3, 0.35, 12, 1, "EIOS — Enterprise Intelligence",
              font_size=14, bold=False, color=_WHITE, align=PP_ALIGN.LEFT)
    _text_box(slide, 0.5, 2.0, 12, 1.5, report_title,
              font_size=36, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)
    _text_box(slide, 0.5, 3.6, 12, 0.5, organization_name,
              font_size=20, bold=False, color=RGBColor(0x93, 0xC5, 0xFD), align=PP_ALIGN.CENTER)
    _text_box(slide, 0.5, 4.2, 12, 0.4, f"{period_label}  ·  {date_str}",
              font_size=12, bold=False, color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.CENTER)

    # ── Slide 2: Executive Summary ────────────────────────────────────────────
    slide = _add_slide(prs)
    _fill_background(slide, _LIGHT)
    _text_box(slide, 0.5, 0.3, 12, 0.6, "Executive Summary",
              font_size=24, bold=True, color=_DARK)
    summary_text = executive_summary[:1200] if executive_summary else "No summary provided."
    _text_box(slide, 0.5, 1.1, 12, 5.5, summary_text, font_size=14, color=_DARK)

    # ── Slide 3: KPI Highlights ───────────────────────────────────────────────
    if kpi_highlights:
        slide = _add_slide(prs)
        _fill_background(slide, _WHITE)
        _text_box(slide, 0.5, 0.2, 12, 0.6, "KPI Highlights",
                  font_size=22, bold=True, color=_DARK)

        cols = min(len(kpi_highlights), 4)
        col_width = 12 / cols
        for i, kpi in enumerate(kpi_highlights[:8]):
            col = i % cols
            row = i // cols
            left = 0.5 + col * col_width
            top = 1.1 + row * 2.5
            shape = slide.shapes.add_shape(
                1, Inches(left), Inches(top), Inches(col_width - 0.2), Inches(2.2)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0xEF, 0xF6, 0xFF)
            _text_box(slide, left + 0.1, top + 0.1, col_width - 0.4, 0.4,
                      kpi.get("label", ""), font_size=10, color=RGBColor(0x64, 0x74, 0x8B))
            _text_box(slide, left + 0.1, top + 0.5, col_width - 0.4, 0.8,
                      f"{kpi.get('value', '—')} {kpi.get('unit', '')}".strip(),
                      font_size=20, bold=True, color=_BLUE)
            trend = kpi.get("trend", "")
            if trend:
                _text_box(slide, left + 0.1, top + 1.4, col_width - 0.4, 0.4,
                          trend, font_size=10,
                          color=_ACCENT if "↑" in trend else RGBColor(0xEF, 0x44, 0x44))

    # ── Slide 4: Risk Register ────────────────────────────────────────────────
    if risks:
        slide = _add_slide(prs)
        _fill_background(slide, _WHITE)
        _text_box(slide, 0.5, 0.2, 12, 0.6, "Risk Register — Top Open Risks",
                  font_size=22, bold=True, color=_DARK)

        sev_color = {"CRITICAL": RGBColor(0xDC, 0x26, 0x26),
                     "HIGH": RGBColor(0xEA, 0x58, 0x0C),
                     "MEDIUM": RGBColor(0xCA, 0x8A, 0x04),
                     "LOW": RGBColor(0x05, 0x96, 0x69)}

        table = slide.shapes.add_table(
            min(len(risks), 10) + 1, 4,
            Inches(0.5), Inches(1.0), Inches(12.3), Inches(5.8)
        ).table

        for i, header in enumerate(["Risk", "Severity", "Status", "Owner"]):
            cell = table.cell(0, i)
            cell.text = header
            cell.text_frame.paragraphs[0].runs[0].font.bold = True
            cell.text_frame.paragraphs[0].runs[0].font.color.rgb = _WHITE
            _set_cell_bg(cell, _BLUE)

        for r_idx, risk in enumerate(risks[:10], 1):
            sev = risk.get("severity", "MEDIUM")
            table.cell(r_idx, 0).text = risk.get("title", "")[:80]
            sc = table.cell(r_idx, 1)
            sc.text = sev
            _set_cell_bg(sc, sev_color.get(sev, RGBColor(0xCA, 0x8A, 0x04)))
            table.cell(r_idx, 2).text = risk.get("status", "")
            table.cell(r_idx, 3).text = risk.get("owner", "")[:40]

    # ── Slide 5: Next Steps ───────────────────────────────────────────────────
    if recommendations:
        slide = _add_slide(prs)
        _fill_background(slide, _LIGHT)
        _text_box(slide, 0.5, 0.2, 12, 0.6, "Next Steps & Recommendations",
                  font_size=22, bold=True, color=_DARK)

        for i, rec in enumerate(recommendations[:8]):
            top = 1.0 + i * 0.75
            prio = rec.get("priority", "Medium")
            color = _DARK if prio == "Low" else (
                RGBColor(0xCA, 0x8A, 0x04) if prio == "Medium" else
                RGBColor(0xDC, 0x26, 0x26)
            )
            _text_box(slide, 0.5, top, 10, 0.6,
                      f"{'●'} [{prio}] {rec.get('title', '')} — due {rec.get('due_date', 'TBD')}",
                      font_size=12, color=color)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
